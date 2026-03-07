"""
NeuralDoc — Document Loader Module
Handles loading and parsing of multiple document formats.
"""

import logging
from pathlib import Path
from datetime import datetime
from typing import Optional

from langchain_core.documents import Document

logger = logging.getLogger(__name__)


class DocumentLoader:
    """Loads documents from various file formats and returns LangChain Document objects."""

    SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".pptx", ".xlsx", ".xls", ".csv", ".md"}

    def __init__(self) -> None:
        self._loaders = {
            ".pdf": self._load_pdf,
            ".docx": self._load_docx,
            ".txt": self._load_txt,
            ".pptx": self._load_pptx,
            ".xlsx": self._load_excel,
            ".xls": self._load_excel,
            ".csv": self._load_csv,
            ".md": self._load_markdown,
        }

    async def load(
        self,
        file_path: str,
        session_id: str,
        original_filename: Optional[str] = None,
    ) -> list[Document]:
        """
        Load a document from the given file path.

        Args:
            file_path: Path to the file on disk.
            session_id: Current chat session ID.
            original_filename: Original name of the uploaded file.

        Returns:
            List of LangChain Document objects with metadata.

        Raises:
            ValueError: If the file type is not supported.
            FileNotFoundError: If the file does not exist.
        """
        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = path.suffix.lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Unsupported file type: '{ext}'. "
                f"Supported: {', '.join(sorted(self.SUPPORTED_EXTENSIONS))}"
            )

        filename = original_filename or path.name
        logger.info("Loading document: %s (type: %s)", filename, ext)

        try:
            documents = self._loaders[ext](str(path))
        except Exception as e:
            logger.error("Failed to load '%s': %s", filename, e)
            raise RuntimeError(f"Failed to load '{filename}': {e}") from e

        # Enrich metadata on every chunk
        base_metadata = {
            "source": filename,
            "file_type": ext,
            "upload_timestamp": datetime.now().isoformat(),
            "session_id": session_id,
        }
        for doc in documents:
            doc.metadata.update(base_metadata)

        logger.info("Loaded %d pages/sections from '%s'", len(documents), filename)
        return documents

    # ── Individual Loaders ─────────────────────────────────────────────────

    @staticmethod
    def _load_pdf(file_path: str) -> list[Document]:
        """Load PDF using PyPDFLoader."""
        from langchain_community.document_loaders import PyPDFLoader
        return PyPDFLoader(file_path).load()

    @staticmethod
    def _load_docx(file_path: str) -> list[Document]:
        """Load DOCX using Docx2txtLoader."""
        from langchain_community.document_loaders import Docx2txtLoader
        return Docx2txtLoader(file_path).load()

    @staticmethod
    def _load_txt(file_path: str) -> list[Document]:
        """Load plain text files."""
        from langchain_community.document_loaders import TextLoader
        return TextLoader(file_path, encoding="utf-8").load()

    @staticmethod
    def _load_pptx(file_path: str) -> list[Document]:
        """Load PowerPoint files using python-pptx."""
        from pptx import Presentation

        prs = Presentation(file_path)
        documents: list[Document] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_parts.append(shape.text_frame.text)
            if text_parts:
                content = "\n".join(text_parts)
                documents.append(
                    Document(
                        page_content=content,
                        metadata={"slide_number": slide_num},
                    )
                )

        return documents

    @staticmethod
    def _load_excel(file_path: str) -> list[Document]:
        """Load Excel files using pandas."""
        import pandas as pd

        documents: list[Document] = []
        xls = pd.ExcelFile(file_path)

        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            if df.empty:
                continue
            content = df.to_string(index=False)
            documents.append(
                Document(
                    page_content=content,
                    metadata={"sheet_name": sheet_name},
                )
            )

        return documents

    @staticmethod
    def _load_csv(file_path: str) -> list[Document]:
        """Load CSV files using pandas."""
        import pandas as pd

        df = pd.read_csv(file_path)
        if df.empty:
            return []
        return [Document(page_content=df.to_string(index=False), metadata={})]

    @staticmethod
    def _load_markdown(file_path: str) -> list[Document]:
        """Load Markdown files."""
        from langchain_community.document_loaders import TextLoader
        return TextLoader(file_path, encoding="utf-8").load()
